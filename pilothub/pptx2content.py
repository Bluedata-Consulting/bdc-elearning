from pptx import Presentation


class PPTxFile:
    def __init__(self, file_path):
        self.prs = Presentation(file_path)
        self.slides = self.prs.slides
    
    def get_slide_text(self, slide):
        """Returns the text in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
        Returns:
            str -- text in slide
        """
        slide_text = '\n'.join([shape.text for shape in slide.shapes if 
                                hasattr(shape, "text")])
        return slide_text

    def get_slide_notes(self, slide):
        """Returns the notes in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
        Returns:
            str -- notes in slide
        """
        notes_slide = self.slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text
        return notes_text

    def set_slide_notes(self, slide, text):
        """Sets the notes in a slide
        
        Arguments:
            slide {pptx.slide} -- slide object
            text {str} -- text to be set as notes
        """
        slide.notes_slide.notes_text_frame.text = text

    def erase_slide_notes(self, slide):
        slide.notes_slide.notes_text_frame.clear()

    def set_skip_slides(self, skip_slides_index: list[int] = None,
                        skip_slides_layout: list[str] = None,
                        skip_slide_titles: list[str] = None):
        """
        Set the slides to skip.
        :param skip_slides_index: Index of the slides to skip.
        :param skip_slides_layout: Layout of the slides to skip.
        :param skip_slide_titles: Titles of the slides to skip.
        """
        if skip_slides_index is None:
            skip_slides_index = [1, -1]
        if skip_slides_layout is None:
            skip_slides_layout = ["Title Slide", "CoverPage", "Quote Slide", 
                                  "Agenda", "Section Header", "QuoteHead"]
        if skip_slide_titles is None:
            skip_slide_titles = []
        self.skip_slides_index = skip_slides_index
        self.skip_slides_layout = skip_slides_layout
        self.skip_slide_titles = skip_slide_titles
